from __future__ import annotations

import json
from pathlib import Path

import qrcode
from PIL import Image, ImageDraw, ImageFont
from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import PP_ALIGN
from pptx.util import Inches, Pt


ROOT = Path("/Users/ghantadivakar/reputation-youtube-monitor")
PUBLIC = ROOT / "public"
BRAND = ROOT / "brand"
ASSETS = ROOT / "scripts" / "presentation_assets_v2"
OUTPUT = ROOT / "output"

DASHBOARD_LINK = "http://localhost:3000/client-demo.html"
ANALYST_LINK = "http://localhost:3000/youtube-dashboard.html"
SHEETS_LINK = "https://docs.google.com/spreadsheets/d/1-hFwvEvEYa0w-iZ808mAt-JpiDGFaqQhGs2rd0SxJvk"

INK = RGBColor(8, 14, 32)
NAVY = RGBColor(10, 22, 50)
BLUE = RGBColor(37, 99, 235)
CYAN = RGBColor(6, 182, 212)
GREEN = RGBColor(16, 185, 129)
RED = RGBColor(220, 38, 38)
AMBER = RGBColor(217, 119, 6)
WHITE = RGBColor(255, 255, 255)
FOG = RGBColor(248, 250, 252)
LINE = RGBColor(226, 232, 240)
MUTED = RGBColor(100, 116, 139)
TEXT = RGBColor(15, 23, 42)


def font_path(bold=False):
    candidates = [
        "/System/Library/Fonts/Supplemental/Arial Bold.ttf" if bold else "/System/Library/Fonts/Supplemental/Arial.ttf",
        "/System/Library/Fonts/Supplemental/Helvetica Bold.ttf" if bold else "/System/Library/Fonts/Supplemental/Helvetica.ttf",
    ]
    for c in candidates:
        if Path(c).exists():
            return c
    return None


def pil_font(size, bold=False):
    p = font_path(bold)
    return ImageFont.truetype(p, size) if p else ImageFont.load_default()


def rgb_tuple(c: RGBColor):
    return (c[0], c[1], c[2])


def add_bg(slide, color=FOG):
    shp = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, Inches(13.333), Inches(7.5))
    shp.fill.solid()
    shp.fill.fore_color.rgb = color
    shp.line.fill.background()


def add_text(slide, text, x, y, w, h, size=18, bold=False, color=TEXT, align=None):
    box = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    tf = box.text_frame
    tf.clear()
    tf.margin_left = 0
    tf.margin_right = 0
    tf.margin_top = 0
    tf.margin_bottom = 0
    p = tf.paragraphs[0]
    p.text = text
    p.font.name = "Arial"
    p.font.size = Pt(size)
    p.font.bold = bold
    p.font.color.rgb = color
    if align:
        p.alignment = align
    return box


def add_link(slide, text, url, x, y, w, h, color=BLUE):
    box = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    tf = box.text_frame
    tf.clear()
    p = tf.paragraphs[0]
    run = p.add_run()
    run.text = text
    run.font.name = "Arial"
    run.font.size = Pt(11)
    run.font.bold = True
    run.font.color.rgb = color
    run.hyperlink.address = url
    return box


def add_rule(slide, x, y, w, color=CYAN):
    shp = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(x), Inches(y), Inches(w), Inches(0.03))
    shp.fill.solid()
    shp.fill.fore_color.rgb = color
    shp.line.fill.background()


def add_kicker(slide, text, x=0.72, y=0.52, color=BLUE):
    add_text(slide, text.upper(), x, y, 3.2, 0.18, 8, True, color)


def add_title(slide, kicker, title, subtitle=None, dark=False):
    add_kicker(slide, kicker, color=CYAN if dark else BLUE)
    add_text(slide, title, 0.72, 0.82, 8.9, 0.62, 26, True, WHITE if dark else TEXT)
    if subtitle:
        add_text(slide, subtitle, 0.74, 1.47, 9.2, 0.38, 11, False, RGBColor(203, 213, 225) if dark else MUTED)
    add_rule(slide, 0.72, 1.98, 1.1, CYAN if dark else BLUE)


def add_footer(slide, number, dark=False):
    c = RGBColor(148, 163, 184) if dark else MUTED
    add_text(slide, "APP.AI", 0.72, 7.1, 1.0, 0.15, 8, True, c)
    add_text(slide, f"{number:02d}", 12.15, 7.1, 0.45, 0.15, 8, True, c, PP_ALIGN.RIGHT)


def add_card(slide, title, body, x, y, w, h, accent=BLUE):
    shp = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(x), Inches(y), Inches(w), Inches(h))
    shp.fill.solid()
    shp.fill.fore_color.rgb = WHITE
    shp.line.color.rgb = LINE
    add_rule(slide, x + 0.24, y + 0.28, 0.45, accent)
    add_text(slide, title, x + 0.24, y + 0.55, w - 0.48, 0.28, 13, True, TEXT)
    add_text(slide, body, x + 0.24, y + 0.94, w - 0.48, h - 1.04, 9, False, MUTED)
    return shp


def add_metric(slide, value, label, x, y, w=1.45, accent=BLUE):
    add_text(slide, value, x, y, w, 0.42, 25, True, accent, PP_ALIGN.CENTER)
    add_text(slide, label.upper(), x, y + 0.45, w, 0.2, 7, True, MUTED, PP_ALIGN.CENTER)


def add_bullets(slide, items, x, y, w, h, size=12, color=TEXT):
    box = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    tf = box.text_frame
    tf.clear()
    tf.margin_left = 0
    tf.margin_right = 0
    for i, item in enumerate(items):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.text = item
        p.font.name = "Arial"
        p.font.size = Pt(size)
        p.font.color.rgb = color
        p.space_after = Pt(8)
    return box


def build_logo():
    BRAND.mkdir(exist_ok=True)
    img = Image.new("RGBA", (1500, 420), rgb_tuple(INK) + (255,))
    d = ImageDraw.Draw(img)
    d.rounded_rectangle((0, 0, 1499, 419), radius=46, fill=rgb_tuple(INK))
    d.rounded_rectangle((75, 85, 255, 265), radius=42, outline=rgb_tuple(CYAN), width=16)
    d.rectangle((75, 174, 165, 265), fill=rgb_tuple(CYAN))
    d.rectangle((165, 210, 255, 265), fill=rgb_tuple(GREEN))
    d.rectangle((190, 85, 255, 265), fill=rgb_tuple(INK))
    d.text((335, 105), "APP", font=pil_font(118, True), fill=rgb_tuple(WHITE))
    d.ellipse((635, 203, 665, 233), fill=rgb_tuple(GREEN))
    d.text((725, 105), "AI", font=pil_font(118, True), fill=rgb_tuple(WHITE))
    d.rounded_rectangle((905, 100, 930, 250), radius=10, fill=rgb_tuple(CYAN))
    d.text((335, 258), "Applied Prediction Platforms", font=pil_font(28, False), fill=(148, 163, 184))
    img.save(BRAND / "app-ai-logo-v2.png")

    svg = """<svg width="1500" height="420" viewBox="0 0 1500 420" fill="none" xmlns="http://www.w3.org/2000/svg">
<rect width="1500" height="420" rx="46" fill="#080E20"/>
<rect x="75" y="85" width="180" height="180" rx="42" stroke="#06B6D4" stroke-width="16"/>
<path d="M75 174H165V265H75V174Z" fill="#06B6D4"/>
<path d="M165 210H255V265H165V210Z" fill="#10B981"/>
<path d="M190 85H255V265H190V85Z" fill="#080E20"/>
<text x="335" y="216" fill="#FFFFFF" font-family="Arial, Helvetica, sans-serif" font-size="118" font-weight="700">APP</text>
<circle cx="650" cy="218" r="15" fill="#10B981"/>
<text x="725" y="216" fill="#FFFFFF" font-family="Arial, Helvetica, sans-serif" font-size="118" font-weight="700">AI</text>
<rect x="905" y="100" width="25" height="150" rx="10" fill="#06B6D4"/>
<text x="335" y="290" fill="#94A3B8" font-family="Arial, Helvetica, sans-serif" font-size="28">Applied Prediction Platforms</text>
</svg>"""
    (BRAND / "app-ai-logo-v2.svg").write_text(svg)


def qr(url, name):
    q = qrcode.QRCode(box_size=9, border=2)
    q.add_data(url)
    q.make(fit=True)
    img = q.make_image(fill_color=rgb_tuple(INK), back_color=(255, 255, 255)).convert("RGB")
    path = ASSETS / name
    img.save(path)
    return path


def donut(path, values, labels, colors):
    img = Image.new("RGB", (950, 620), rgb_tuple(FOG))
    d = ImageDraw.Draw(img)
    d.rounded_rectangle((20, 20, 930, 600), radius=22, fill=rgb_tuple(WHITE), outline=rgb_tuple(LINE), width=2)
    total = sum(values) or 1
    start = -90
    box = (90, 110, 470, 490)
    for val, col in zip(values, colors):
        end = start + 360 * val / total
        d.pieslice(box, start, end, fill=rgb_tuple(col))
        start = end
    d.ellipse((210, 230, 350, 370), fill=rgb_tuple(WHITE))
    d.text((232, 242), str(total), font=pil_font(42, True), fill=rgb_tuple(TEXT))
    d.text((226, 298), "items", font=pil_font(20, False), fill=rgb_tuple(MUTED))
    y = 155
    for label, val, col in zip(labels, values, colors):
        pct = round(val * 100 / total)
        d.rounded_rectangle((560, y, 590, y + 28), radius=7, fill=rgb_tuple(col))
        d.text((615, y - 1), f"{label}", font=pil_font(26, True), fill=rgb_tuple(TEXT))
        d.text((615, y + 33), f"{val} items | {pct}%", font=pil_font(20, False), fill=rgb_tuple(MUTED))
        y += 90
    img.save(path)
    return path


def bars(path, labels, values, colors, title):
    img = Image.new("RGB", (1150, 620), rgb_tuple(FOG))
    d = ImageDraw.Draw(img)
    d.rounded_rectangle((20, 20, 1130, 600), radius=22, fill=rgb_tuple(WHITE), outline=rgb_tuple(LINE), width=2)
    d.text((60, 55), title, font=pil_font(30, True), fill=rgb_tuple(TEXT))
    max_v = max(values) if values else 1
    base_y = 500
    left = 95
    gap = 55
    bar_w = (930 - gap * (len(values) - 1)) / len(values)
    for i, (label, val, color) in enumerate(zip(labels, values, colors)):
        x = left + i * (bar_w + gap)
        h = int(330 * val / max_v)
        d.rounded_rectangle((x, base_y - h, x + bar_w, base_y), radius=14, fill=rgb_tuple(color))
        d.text((x, base_y - h - 40), str(val), font=pil_font(30, True), fill=rgb_tuple(TEXT))
        d.text((x, base_y + 20), label, font=pil_font(21, True), fill=rgb_tuple(MUTED))
    img.save(path)
    return path


def workflow(path):
    img = Image.new("RGB", (1500, 520), rgb_tuple(FOG))
    d = ImageDraw.Draw(img)
    steps = [
        ("Collect", "YouTube, Meta, X, web"),
        ("Store", "Sheets, DB, APIs"),
        ("Analyze", "AI sentiment + summaries"),
        ("Act", "Dashboard, alerts, follow-up"),
    ]
    x = 70
    for i, (a, b) in enumerate(steps):
        d.rounded_rectangle((x, 125, x + 300, 390), radius=20, fill=rgb_tuple(WHITE), outline=rgb_tuple(LINE), width=2)
        d.text((x + 32, 162), f"0{i+1}", font=pil_font(18, True), fill=rgb_tuple(CYAN))
        d.text((x + 32, 215), a, font=pil_font(36, True), fill=rgb_tuple(TEXT))
        d.text((x + 32, 284), b, font=pil_font(22, False), fill=rgb_tuple(MUTED))
        if i < 3:
            d.line((x + 320, 260, x + 380, 260), fill=rgb_tuple(BLUE), width=5)
            d.polygon([(x + 380, 260), (x + 355, 245), (x + 355, 275)], fill=rgb_tuple(BLUE))
        x += 365
    img.save(path)
    return path


def build_deck():
    BRAND.mkdir(exist_ok=True)
    ASSETS.mkdir(parents=True, exist_ok=True)
    OUTPUT.mkdir(exist_ok=True)
    build_logo()

    data = json.loads((PUBLIC / "youtube_dashboard_data.json").read_text())
    latest = data["latest"]
    total = latest["totalVideos"]
    positive = latest["positive"]
    negative = latest["negative"]
    neutral = latest["neutral"]
    videos = latest["videos"]
    shorts = latest["shorts"]
    risk = latest["riskScore"]
    date = latest["date"]

    chart_sent = donut(ASSETS / "sentiment.png", [positive, negative, neutral], ["Positive", "Negative", "Neutral"], [GREEN, RED, MUTED])
    chart_formats = bars(ASSETS / "formats.png", ["Videos", "Shorts"], [videos, shorts], [BLUE, CYAN], "Content format split")
    flow = workflow(ASSETS / "workflow.png")
    qr_dash = qr(DASHBOARD_LINK, "qr_dashboard.png")
    qr_sheet = qr(SHEETS_LINK, "qr_sheet.png")

    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)
    blank = prs.slide_layouts[6]

    # 1
    s = prs.slides.add_slide(blank)
    add_bg(s, INK)
    s.shapes.add_picture(str(BRAND / "app-ai-logo-v2.png"), Inches(0.72), Inches(0.52), width=Inches(3.3))
    add_text(s, "AI Intelligence,\nAutomation and Data Systems", 0.82, 1.95, 8.8, 1.15, 30, True, WHITE)
    add_text(s, "For reputation monitoring, data analysis, ML models, AI agents and multilingual business automation.", 0.84, 3.25, 8.0, 0.38, 14, False, RGBColor(203, 213, 225))
    add_rule(s, 0.84, 3.95, 1.4, CYAN)
    add_text(s, "Client capability deck", 0.84, 4.2, 3.0, 0.25, 10, True, CYAN)
    add_footer(s, 1, True)

    # 2
    s = prs.slides.add_slide(blank)
    add_bg(s)
    add_title(s, "Company Overview", "APP.AI helps clients convert scattered data into clear business decisions.")
    add_card(s, "Positioning", "An AI intelligence and automation company for public figures, agencies, founders and high-visibility organizations.", 0.8, 2.4, 3.75, 2.0, BLUE)
    add_card(s, "Core capability", "We build dashboards, ML models, AI agents, call automations and data workflows that operate daily.", 4.8, 2.4, 3.75, 2.0, CYAN)
    add_card(s, "Business outcome", "Clients get decision-ready insights, risk visibility, faster follow-up and automated reporting.", 8.8, 2.4, 3.75, 2.0, GREEN)
    add_bullets(s, ["Mission: build practical AI-powered systems that collect, analyze, automate and report.", "Vision: become a trusted AI intelligence layer for reputation, growth and operations."], 1.0, 5.25, 10.8, 0.75, 15)
    add_footer(s, 2)

    # 3
    s = prs.slides.add_slide(blank)
    add_bg(s)
    add_title(s, "What We Build", "Four service lines that connect strategy with working systems.")
    items = [
        ("Data Intelligence", "Dashboards, reports, Sheets, databases, APIs and analytics workflows."),
        ("Machine Learning", "Classification, scoring, prediction and recommendation models."),
        ("AI Agents", "Research, monitoring, sales, support and operations agents."),
        ("Automation Systems", "Calls, follow-ups, CRM updates, alerts and daily reporting."),
    ]
    for i, (t, b) in enumerate(items):
        add_card(s, t, b, 0.8 + i * 3.1, 2.35, 2.75, 2.35, [BLUE, CYAN, GREEN, AMBER][i])
    add_footer(s, 3)

    # 4
    s = prs.slides.add_slide(blank)
    add_bg(s)
    add_title(s, "Flagship Product", "Sentiment and Reputation Intelligence Platform.", "Monitor online reputation for a person, politician, celebrity, movie, brand or organization.")
    add_bullets(s, ["Classifies public links as positive, negative or neutral.", "Highlights risky links with summaries for PR or legal review.", "Creates daily dashboards and Google Sheets proof for every item.", "Can expand from YouTube to Meta, X and the open web through compliant data access."], 0.9, 2.5, 5.2, 2.2, 15)
    add_card(s, "Current working demo", f"Thalapathy Vijay YouTube monitoring, last 24 hours. Latest run: {date}.", 7.0, 2.4, 4.8, 1.3, BLUE)
    add_metric(s, str(total), "items", 7.15, 4.25, 1.25, BLUE)
    add_metric(s, str(positive), "positive", 8.55, 4.25, 1.25, GREEN)
    add_metric(s, str(negative), "negative", 9.95, 4.25, 1.25, RED)
    add_metric(s, f"{risk}/100", "risk", 11.35, 4.25, 1.25, AMBER)
    add_footer(s, 4)

    # 5
    s = prs.slides.add_slide(blank)
    add_bg(s)
    add_title(s, "Demo Analytics", "The dashboard converts raw links into executive-level sentiment.")
    s.shapes.add_picture(str(chart_sent), Inches(0.78), Inches(2.15), width=Inches(5.55))
    s.shapes.add_picture(str(chart_formats), Inches(6.75), Inches(2.15), width=Inches(5.65))
    add_footer(s, 5)

    # 6
    s = prs.slides.add_slide(blank)
    add_bg(s)
    add_title(s, "Operating Model", "A simple repeatable workflow from public conversation to action.")
    s.shapes.add_picture(str(flow), Inches(0.65), Inches(2.25), width=Inches(12.0))
    add_text(s, "The client never has to manually search hundreds of links. The system collects, analyzes, scores and presents what needs attention.", 1.0, 6.1, 10.9, 0.35, 14, True, TEXT, PP_ALIGN.CENTER)
    add_footer(s, 6)

    # 7
    s = prs.slides.add_slide(blank)
    add_bg(s)
    add_title(s, "Channel Coverage", "The platform can expand into a full public-conversation intelligence system.")
    channels = [
        ("YouTube", "Videos, Shorts, titles, descriptions, transcripts where available."),
        ("Meta", "Instagram, Facebook, reels, posts and public page data through compliant access."),
        ("X", "Posts, hashtags, mentions and public conversation tracking."),
        ("Web", "News, blogs, websites and public pages with summaries."),
    ]
    for i, (t, b) in enumerate(channels):
        add_card(s, t, b, 0.8 + (i % 2) * 6.1, 2.25 + (i // 2) * 1.8, 5.45, 1.35, [BLUE, GREEN, INK, AMBER][i])
    add_footer(s, 7)

    # 8
    s = prs.slides.add_slide(blank)
    add_bg(s)
    add_title(s, "AI Call Agents", "Inbound and outbound automation in regional languages.")
    add_card(s, "Inbound AI calls", "Answer FAQs, capture leads, book appointments, route urgent calls and update Sheets or CRM.", 0.9, 2.3, 3.75, 2.15, BLUE)
    add_card(s, "Outbound AI calls", "Lead follow-up, reminders, surveys, campaign calls, sales qualification and feedback collection.", 4.85, 2.3, 3.75, 2.15, GREEN)
    add_card(s, "Languages", "Telugu, Tamil, Hindi, Kannada, Marathi, English and more based on client market.", 8.8, 2.3, 3.45, 2.15, CYAN)
    add_text(s, "Best fit: political campaigns, agencies, clinics, education, real estate, local services and high-volume sales teams.", 1.0, 5.45, 10.9, 0.35, 15, True, TEXT, PP_ALIGN.CENTER)
    add_footer(s, 8)

    # 9
    s = prs.slides.add_slide(blank)
    add_bg(s)
    add_title(s, "Target Clients", "Best customers already spend on reputation, campaigns, sales or support.")
    add_bullets(s, ["PR agencies and celebrity management firms", "Political consultants and campaign teams", "Movie production houses and entertainment marketers", "Digital marketing agencies and business founders", "Real estate, education, healthcare and high-visibility local brands"], 0.95, 2.35, 5.4, 2.6, 16)
    add_card(s, "Decision makers", "Founder, co-founder, managing director, PR head, campaign manager, digital marketing head, celebrity manager, political consultant and business development head.", 7.1, 2.35, 4.8, 2.0, BLUE)
    add_card(s, "First market focus", "Andhra Pradesh, Telangana, Karnataka and Mumbai.", 7.1, 4.75, 4.8, 1.05, GREEN)
    add_footer(s, 9)

    # 10
    s = prs.slides.add_slide(blank)
    add_bg(s)
    add_title(s, "Why APP.AI", "Practical AI systems, built around the client’s real workflow.")
    reasons = [
        ("Working proof", f"Live dashboard already processes {total} public YouTube items in the demo run."),
        ("Actionable output", "Links, summaries, scores and next steps instead of only high-level charts."),
        ("Regional strength", "Designed for Indian markets, languages, PR realities and campaign workflows."),
        ("Automation depth", "Analytics, agents, Sheets, CRM, calls and alerts can work together."),
    ]
    for i, (t, b) in enumerate(reasons):
        add_card(s, t, b, 0.8 + (i % 2) * 6.05, 2.2 + (i // 2) * 1.75, 5.45, 1.25, [BLUE, GREEN, CYAN, AMBER][i])
    add_footer(s, 10)

    # 11
    s = prs.slides.add_slide(blank)
    add_bg(s, INK)
    add_title(s, "Demo Access", "Use these links during a client conversation.", "Open dashboard first, then use Google Sheets as raw-data proof.", True)
    s.shapes.add_picture(str(qr_dash), Inches(1.2), Inches(2.55), width=Inches(1.5))
    add_text(s, "Client Dashboard", 2.95, 2.62, 3.0, 0.25, 15, True, WHITE)
    add_link(s, DASHBOARD_LINK, DASHBOARD_LINK, 2.95, 3.02, 5.0, 0.22, CYAN)
    s.shapes.add_picture(str(qr_sheet), Inches(1.2), Inches(4.55), width=Inches(1.5))
    add_text(s, "Google Sheets Data", 2.95, 4.62, 3.0, 0.25, 15, True, WHITE)
    add_link(s, SHEETS_LINK, SHEETS_LINK, 2.95, 5.02, 7.4, 0.22, GREEN)
    add_text(s, "Note: localhost links work on this laptop. For external clients, deploy the dashboard online and replace the URLs.", 1.0, 6.38, 10.9, 0.25, 10, False, RGBColor(203, 213, 225))
    add_footer(s, 11, True)

    # 12
    s = prs.slides.add_slide(blank)
    add_bg(s, INK)
    s.shapes.add_picture(str(BRAND / "app-ai-logo-v2.png"), Inches(0.78), Inches(0.68), width=Inches(3.2))
    add_text(s, "Build intelligence. Automate action. Make faster decisions.", 0.88, 2.6, 8.9, 0.95, 34, True, WHITE)
    add_text(s, "APP.AI helps clients move from data overload to AI-powered decisions through analytics, agents, automation and intelligent dashboards.", 0.9, 3.85, 8.4, 0.45, 14, False, RGBColor(203, 213, 225))
    add_rule(s, 0.9, 4.62, 1.6, CYAN)
    add_footer(s, 12, True)

    out = OUTPUT / "APP_AI_Client_Presentation_Professional.pptx"
    prs.save(out)
    return out


if __name__ == "__main__":
    print(build_deck())
