from __future__ import annotations

import json
import math
from pathlib import Path
from typing import Iterable

import qrcode
from PIL import Image, ImageDraw, ImageFont
from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import PP_ALIGN
from pptx.util import Inches, Pt


ROOT = Path("/Users/ghantadivakar/reputation-youtube-monitor")
PUBLIC = ROOT / "public"
BRAND = ROOT / "brand"
ASSETS = ROOT / "scripts" / "presentation_assets"
OUTPUT = ROOT / "output"

DASHBOARD_LINK = "http://localhost:3000/client-demo.html"
ANALYST_LINK = "http://localhost:3000/youtube-dashboard.html"
SHEETS_LINK = "https://docs.google.com/spreadsheets/d/1-hFwvEvEYa0w-iZ808mAt-JpiDGFaqQhGs2rd0SxJvk"

INK = RGBColor(11, 16, 32)
NAVY = RGBColor(15, 23, 42)
BLUE = RGBColor(29, 78, 216)
CYAN = RGBColor(0, 194, 255)
GREEN = RGBColor(22, 163, 74)
MINT = RGBColor(119, 246, 199)
SLATE = RGBColor(100, 116, 139)
LIGHT = RGBColor(248, 250, 252)
PANEL = RGBColor(241, 245, 249)
RED = RGBColor(239, 68, 68)
AMBER = RGBColor(245, 158, 11)


def font(size: int, bold: bool = False, color: RGBColor = INK):
    return {"name": "Arial", "size": Pt(size), "bold": bold, "color": color}


def set_text(tf, text: str, size: int = 20, bold: bool = False, color: RGBColor = INK):
    tf.clear()
    p = tf.paragraphs[0]
    p.text = text
    p.font.name = "Arial"
    p.font.size = Pt(size)
    p.font.bold = bold
    p.font.color.rgb = color
    return p


def add_textbox(slide, text, x, y, w, h, size=20, bold=False, color=INK, align=None):
    shape = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    p = set_text(shape.text_frame, text, size, bold, color)
    if align:
        p.alignment = align
    return shape


def add_label(slide, text, x, y, w, h, fill=RGBColor(232, 240, 255), color=BLUE):
    shp = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(x), Inches(y), Inches(w), Inches(h))
    shp.fill.solid()
    shp.fill.fore_color.rgb = fill
    shp.line.color.rgb = fill
    set_text(shp.text_frame, text, 10, True, color)
    shp.text_frame.paragraphs[0].alignment = PP_ALIGN.CENTER
    return shp


def add_link(slide, label, url, x, y, w, h, color=CYAN):
    shape = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    tf = shape.text_frame
    tf.clear()
    p = tf.paragraphs[0]
    run = p.add_run()
    run.text = label
    run.font.name = "Arial"
    run.font.size = Pt(12)
    run.font.bold = True
    run.font.color.rgb = color
    run.hyperlink.address = url
    return shape


def add_footer(slide, idx: int):
    add_textbox(slide, "APP.AI", 0.45, 7.15, 1.1, 0.18, 8, True, SLATE)
    add_textbox(slide, f"{idx:02d}", 12.35, 7.15, 0.45, 0.18, 8, True, SLATE, PP_ALIGN.RIGHT)


def add_title(slide, kicker: str, title: str, subtitle: str | None = None):
    add_label(slide, kicker.upper(), 0.65, 0.45, 1.8, 0.32)
    add_textbox(slide, title, 0.65, 0.92, 8.2, 0.62, 27, True, INK)
    if subtitle:
        add_textbox(slide, subtitle, 0.68, 1.55, 9.8, 0.42, 12, False, SLATE)


def add_card(slide, x, y, w, h, title, body, accent=BLUE):
    shp = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(x), Inches(y), Inches(w), Inches(h))
    shp.fill.solid()
    shp.fill.fore_color.rgb = RGBColor(255, 255, 255)
    shp.line.color.rgb = RGBColor(226, 232, 240)
    bar = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(x), Inches(y), Inches(0.06), Inches(h))
    bar.fill.solid()
    bar.fill.fore_color.rgb = accent
    bar.line.color.rgb = accent
    add_textbox(slide, title, x + 0.2, y + 0.18, w - 0.35, 0.3, 14, True, INK)
    add_textbox(slide, body, x + 0.2, y + 0.58, w - 0.35, h - 0.7, 10, False, SLATE)
    return shp


def find_font(candidates: Iterable[str]) -> str | None:
    for c in candidates:
        p = Path(c)
        if p.exists():
            return str(p)
    return None


FONT_BOLD = find_font(
    [
        "/System/Library/Fonts/Supplemental/Arial Bold.ttf",
        "/System/Library/Fonts/Supplemental/Helvetica Bold.ttf",
        "/System/Library/Fonts/SFNS.ttf",
    ]
)
FONT_REG = find_font(
    [
        "/System/Library/Fonts/Supplemental/Arial.ttf",
        "/System/Library/Fonts/Supplemental/Helvetica.ttf",
        "/System/Library/Fonts/SFNS.ttf",
    ]
)


def load_font(size: int, bold: bool = False):
    path = FONT_BOLD if bold else FONT_REG
    return ImageFont.truetype(path, size) if path else ImageFont.load_default()


def create_logo_assets():
    BRAND.mkdir(exist_ok=True)
    ASSETS.mkdir(parents=True, exist_ok=True)

    svg = """<svg width="1200" height="360" viewBox="0 0 1200 360" fill="none" xmlns="http://www.w3.org/2000/svg">
<rect width="1200" height="360" rx="58" fill="#0B1020"/>
<path d="M107 222C107 156 160 103 226 103H284C333 103 373 143 373 192V253H321V196C321 175 304 158 283 158H230C191 158 159 190 159 229V253H107V222Z" fill="#00C2FF"/>
<path d="M190 207H292V253H190V207Z" fill="#77F6C7"/>
<path d="M422 103H558C611 103 654 146 654 199C654 252 611 295 558 295H493V253H556C586 253 610 229 610 199C610 169 586 145 556 145H474V295H422V103Z" fill="#F8FAFC"/>
<path d="M686 103H822C875 103 918 146 918 199C918 252 875 295 822 295H757V253H820C850 253 874 229 874 199C874 169 850 145 820 145H738V295H686V103Z" fill="#F8FAFC"/>
<circle cx="964" cy="270" r="25" fill="#16A34A"/>
<path d="M1015 103H1068L1138 295H1084L1072 258H1009L997 295H944L1015 103ZM1022 218H1059L1041 160L1022 218Z" fill="#F8FAFC"/>
<path d="M1155 103H1207V295H1155V103Z" fill="#00C2FF"/>
<path d="M91 75H1110" stroke="#1D4ED8" stroke-width="6" stroke-linecap="round" opacity=".45"/>
</svg>"""
    (BRAND / "app-ai-logo.svg").write_text(svg)

    img = Image.new("RGBA", (1600, 480), (11, 16, 32, 255))
    d = ImageDraw.Draw(img)
    d.rounded_rectangle((0, 0, 1599, 479), radius=72, fill=(11, 16, 32, 255))
    d.line((115, 100, 1480, 100), fill=(29, 78, 216, 120), width=8)
    d.pieslice((130, 135, 360, 365), 90, 360, fill=(0, 194, 255))
    d.rectangle((245, 275, 392, 340), fill=(119, 246, 199))
    d.rounded_rectangle((300, 135, 460, 365), radius=34, fill=(11, 16, 32, 255))
    title_font = load_font(184, True)
    d.text((545, 134), "APP", font=title_font, fill=(248, 250, 252))
    d.ellipse((1016, 318, 1078, 380), fill=(22, 163, 74))
    d.text((1110, 134), "AI", font=title_font, fill=(248, 250, 252))
    d.rectangle((1396, 134, 1462, 365), fill=(0, 194, 255))
    img.save(BRAND / "app-ai-logo.png")

    icon = Image.new("RGBA", (512, 512), (11, 16, 32, 255))
    d = ImageDraw.Draw(icon)
    d.rounded_rectangle((20, 20, 492, 492), radius=96, fill=(11, 16, 32, 255), outline=(29, 78, 216), width=8)
    d.pieslice((105, 105, 345, 345), 90, 360, fill=(0, 194, 255))
    d.rectangle((225, 250, 380, 325), fill=(119, 246, 199))
    d.rounded_rectangle((288, 105, 420, 345), radius=38, fill=(11, 16, 32, 255))
    d.ellipse((366, 360, 420, 414), fill=(22, 163, 74))
    icon.save(BRAND / "app-ai-icon.png")


def make_qr(url: str, name: str) -> Path:
    qr = qrcode.QRCode(box_size=8, border=2)
    qr.add_data(url)
    qr.make(fit=True)
    img = qr.make_image(fill_color=(11, 16, 32), back_color=(255, 255, 255)).convert("RGB")
    path = ASSETS / name
    img.save(path)
    return path


def make_bar_chart(labels, values, colors, path: Path, title: str | None = None):
    w, h = 1200, 620
    img = Image.new("RGB", (w, h), (248, 250, 252))
    d = ImageDraw.Draw(img)
    d.rounded_rectangle((20, 20, w - 20, h - 20), radius=28, fill=(255, 255, 255), outline=(226, 232, 240), width=2)
    if title:
        d.text((60, 48), title, font=load_font(34, True), fill=(11, 16, 32))
    max_v = max(values) if values else 1
    x0, y0 = 90, 500
    usable_w = 980
    gap = 58
    bar_w = (usable_w - gap * (len(values) - 1)) / max(1, len(values))
    for i, (lab, val, col) in enumerate(zip(labels, values, colors)):
        x = x0 + i * (bar_w + gap)
        bh = int(350 * val / max_v)
        d.rounded_rectangle((x, y0 - bh, x + bar_w, y0), radius=16, fill=col)
        d.text((x, y0 - bh - 42), str(val), font=load_font(32, True), fill=(11, 16, 32))
        d.text((x, y0 + 18), lab, font=load_font(25, True), fill=(100, 116, 139))
    img.save(path)
    return path


def make_donut(values, colors, labels, path: Path):
    img = Image.new("RGB", (900, 620), (248, 250, 252))
    d = ImageDraw.Draw(img)
    d.rounded_rectangle((20, 20, 880, 600), radius=28, fill=(255, 255, 255), outline=(226, 232, 240), width=2)
    box = (85, 95, 485, 495)
    total = sum(values) or 1
    start = -90
    for val, col in zip(values, colors):
        end = start + 360 * val / total
        d.pieslice(box, start, end, fill=col)
        start = end
    d.ellipse((205, 215, 365, 375), fill=(255, 255, 255))
    d.text((240, 244), str(total), font=load_font(44, True), fill=(11, 16, 32))
    d.text((224, 300), "videos", font=load_font(22, False), fill=(100, 116, 139))
    y = 155
    for lab, val, col in zip(labels, values, colors):
        pct = round(val / total * 100)
        d.rounded_rectangle((560, y, 590, y + 30), radius=8, fill=col)
        d.text((610, y - 2), f"{lab}: {val} ({pct}%)", font=load_font(27, True), fill=(11, 16, 32))
        y += 70
    img.save(path)
    return path


def make_pipeline(path: Path):
    img = Image.new("RGB", (1500, 530), (248, 250, 252))
    d = ImageDraw.Draw(img)
    steps = [
        ("Collect", "YouTube, Meta, X, web"),
        ("Understand", "transcripts, captions, descriptions"),
        ("Score", "positive, negative, neutral, risk"),
        ("Act", "alerts, sheets, dashboard, PR actions"),
    ]
    colors = [(29, 78, 216), (0, 194, 255), (22, 163, 74), (245, 158, 11)]
    x = 75
    for i, (title, body) in enumerate(steps):
        d.rounded_rectangle((x, 115, x + 285, 405), radius=26, fill=(255, 255, 255), outline=(226, 232, 240), width=2)
        d.ellipse((x + 32, 145, x + 92, 205), fill=colors[i])
        d.text((x + 51, 158), str(i + 1), font=load_font(28, True), fill=(255, 255, 255))
        d.text((x + 32, 235), title, font=load_font(34, True), fill=(11, 16, 32))
        d.text((x + 32, 292), body, font=load_font(22, False), fill=(100, 116, 139))
        if i < len(steps) - 1:
            d.line((x + 305, 260, x + 375, 260), fill=(29, 78, 216), width=5)
            d.polygon([(x + 375, 260), (x + 350, 245), (x + 350, 275)], fill=(29, 78, 216))
        x += 360
    img.save(path)
    return path


def bg(slide, color=LIGHT):
    rect = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, Inches(13.333), Inches(7.5))
    rect.fill.solid()
    rect.fill.fore_color.rgb = color
    rect.line.fill.background()


def add_bullets(slide, items, x, y, w, h, size=13, color=INK):
    box = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    tf = box.text_frame
    tf.clear()
    for i, item in enumerate(items):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.text = item
        p.font.name = "Arial"
        p.font.size = Pt(size)
        p.font.color.rgb = color
        p.level = 0
        p.space_after = Pt(8)
    return box


def add_metric(slide, x, y, value, label, accent=BLUE):
    add_textbox(slide, value, x, y, 1.9, 0.44, 27, True, accent)
    add_textbox(slide, label, x, y + 0.48, 2.1, 0.28, 9, True, SLATE)


def build_deck():
    create_logo_assets()
    ASSETS.mkdir(parents=True, exist_ok=True)
    OUTPUT.mkdir(exist_ok=True)

    data = json.loads((PUBLIC / "youtube_dashboard_data.json").read_text())
    latest = data["latest"]
    total = latest["totalVideos"]
    positive = latest["positive"]
    negative = latest["negative"]
    neutral = latest["neutral"]
    risk = latest["riskScore"]
    date = latest["date"]
    videos = latest["videos"]
    shorts = latest["shorts"]
    top_channels = latest["topChannels"][:5]

    qr_dashboard = make_qr(DASHBOARD_LINK, "qr_client_dashboard.png")
    qr_sheet = make_qr(SHEETS_LINK, "qr_google_sheet.png")
    donut = make_donut(
        [positive, negative, neutral],
        [(22, 163, 74), (239, 68, 68), (100, 116, 139)],
        ["Positive", "Negative", "Neutral"],
        ASSETS / "sentiment_donut.png",
    )
    format_chart = make_bar_chart(
        ["Videos", "Shorts"],
        [videos, shorts],
        [(29, 78, 216), (0, 194, 255)],
        ASSETS / "format_chart.png",
        "Content format split",
    )
    channel_chart = make_bar_chart(
        [c["channel"][:14] for c in top_channels],
        [c["count"] for c in top_channels],
        [(29, 78, 216), (0, 194, 255), (22, 163, 74), (245, 158, 11), (100, 116, 139)],
        ASSETS / "channel_chart.png",
        "Top channels in demo data",
    )
    pipeline = make_pipeline(ASSETS / "pipeline.png")

    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)

    def blank():
        return prs.slides.add_slide(prs.slide_layouts[6])

    # 1 cover
    s = blank()
    bg(s, INK)
    s.shapes.add_picture(str(BRAND / "app-ai-logo.png"), Inches(0.72), Inches(0.58), width=Inches(3.1))
    add_textbox(s, "AI systems for reputation, growth and operations", 0.82, 2.15, 8.8, 0.72, 32, True, LIGHT)
    add_textbox(
        s,
        "Data analysis | ML models | AI agents | multilingual automations",
        0.85,
        3.05,
        7.1,
        0.35,
        15,
        False,
        RGBColor(203, 213, 225),
    )
    add_label(s, "CLIENT DEMO DECK", 0.86, 3.72, 2.05, 0.36, RGBColor(17, 36, 67), CYAN)
    for x, y, c in [(10.1, 1.15, CYAN), (11.05, 2.35, GREEN), (9.25, 4.55, BLUE), (11.65, 5.15, AMBER)]:
        shp = s.shapes.add_shape(MSO_SHAPE.OVAL, Inches(x), Inches(y), Inches(0.45), Inches(0.45))
        shp.fill.solid()
        shp.fill.fore_color.rgb = c
        shp.line.fill.background()
    add_textbox(s, "Prepared for client conversations", 0.85, 6.78, 4.6, 0.25, 10, False, RGBColor(148, 163, 184))
    add_footer(s, 1)

    # 2 mission
    s = blank()
    bg(s)
    add_title(s, "Company", "APP.AI turns scattered business data into daily decisions.")
    add_card(s, 0.8, 2.0, 5.6, 1.55, "Mission", "Build practical AI agents that help teams see what is happening, understand why it matters and act faster with confidence.", BLUE)
    add_card(s, 6.85, 2.0, 5.6, 1.55, "Vision", "Become the intelligence layer for brands, public figures, campaigns and growing companies across India and global markets.", CYAN)
    add_card(s, 0.8, 4.1, 3.65, 1.5, "Data analysis", "Dashboards, daily reports, Sheets workflows and executive summaries.", GREEN)
    add_card(s, 4.85, 4.1, 3.65, 1.5, "ML models", "Prediction, classification, scoring and recommendations tailored to client data.", BLUE)
    add_card(s, 8.9, 4.1, 3.65, 1.5, "AI agents", "Agents for research, monitoring, calls, sales follow-up and operations.", AMBER)
    add_footer(s, 2)

    # 3 what we build
    s = blank()
    bg(s)
    add_title(s, "Platform", "Four capability lines, one execution partner.", "APP.AI is positioned as a build-and-operate AI company, not only a dashboard vendor.")
    cards = [
        ("Reputation Intelligence", "Track public sentiment across YouTube, Meta, X and the open web. Find risk links and positive amplifiers early.", BLUE),
        ("Business Analytics", "Connect Sheets, CRMs and databases into decision dashboards for founders and teams.", GREEN),
        ("Custom ML Models", "Train or configure models for classification, forecasting, scoring and workflow recommendations.", CYAN),
        ("Automation Agents", "Automate research, lead qualification, inbound calls, outbound campaigns and follow-ups.", AMBER),
    ]
    for i, (t, b, c) in enumerate(cards):
        add_card(s, 0.8 + (i % 2) * 6.1, 2.0 + (i // 2) * 1.85, 5.55, 1.38, t, b, c)
    add_footer(s, 3)

    # 4 demo proof
    s = blank()
    bg(s)
    add_title(s, "Proof", "Live demo: Thalapathy Vijay reputation analysis.", f"Latest automated run: {date}. Data is written to Google Sheets and visualized in the dashboard.")
    add_metric(s, 0.9, 2.25, str(total), "YouTube items in 24 hours", BLUE)
    add_metric(s, 3.15, 2.25, str(positive), "positive", GREEN)
    add_metric(s, 5.05, 2.25, str(negative), "negative", RED)
    add_metric(s, 6.82, 2.25, str(neutral), "neutral", SLATE)
    add_metric(s, 8.62, 2.25, f"{risk}/100", "risk score", AMBER)
    s.shapes.add_picture(str(donut), Inches(0.8), Inches(3.4), width=Inches(5.25))
    s.shapes.add_picture(str(format_chart), Inches(6.35), Inches(3.4), width=Inches(5.7))
    add_footer(s, 4)

    # 5 product workflow
    s = blank()
    bg(s)
    add_title(s, "Workflow", "From public conversation to client-ready action list.", "The agent pipeline is designed to support PR, legal, campaign and growth teams.")
    s.shapes.add_picture(str(pipeline), Inches(0.55), Inches(2.1), width=Inches(12.15))
    add_bullets(
        s,
        [
            "Current demo analyzes YouTube title, description and available transcripts.",
            "Next phase can add deeper transcript extraction, comment sampling, Meta/X/web data providers and alerting.",
            "Every risky or positive item keeps the original link so the client can verify before action.",
        ],
        1.0,
        6.05,
        11.4,
        0.65,
        11,
        SLATE,
    )
    add_footer(s, 5)

    # 6 multi channel
    s = blank()
    bg(s)
    add_title(s, "Sentiment Intelligence", "Coverage expands from YouTube to the full public conversation.")
    add_card(s, 0.8, 2.05, 2.8, 2.0, "YouTube", "Videos, Shorts, titles, descriptions, channels, transcripts where available, links and risk scoring.", BLUE)
    add_card(s, 3.95, 2.05, 2.8, 2.0, "Meta", "Instagram posts, reels and public page content through approved APIs, partner tools or compliant exports.", GREEN)
    add_card(s, 7.1, 2.05, 2.8, 2.0, "X", "Keyword, hashtag, mention and post-level tracking through API or compliant social listening providers.", INK)
    add_card(s, 10.25, 2.05, 2.8, 2.0, "Web", "News, blogs, forums and web pages with link-level summaries and daily trend movement.", AMBER)
    add_bullets(
        s,
        [
            "Outputs: positive/negative/neutral percentages, risky links, top supportive links, summaries and recommended PR actions.",
            "Best practice: use official APIs and approved data providers for platforms with strict scraping rules.",
            "Human-in-loop review is kept for legal, political and brand-sensitive decisions.",
        ],
        0.95,
        4.75,
        11.7,
        1.25,
        14,
        INK,
    )
    add_footer(s, 6)

    # 7 use cases
    s = blank()
    bg(s)
    add_title(s, "Uses", "Who benefits immediately from this product?")
    use_cases = [
        ("Celebrities & managers", "Track reputation, fan sentiment, controversies, movie launch perception and influencer support."),
        ("Politicians & campaigns", "Monitor constituency narrative, opposition attacks, media tone and issue-level risk."),
        ("PR agencies", "Daily client reports, crisis detection, influencer mapping and proof-backed campaign insights."),
        ("Movie production houses", "Pre-release hype, trailer response, actor sentiment, review risk and channel-wise coverage."),
        ("Brands & founders", "Customer voice, product sentiment, competitor comparison and executive escalation."),
        ("Digital agencies", "Add AI monitoring and automation as a premium retainer product."),
    ]
    for i, (t, b) in enumerate(use_cases):
        add_card(s, 0.75 + (i % 3) * 4.2, 1.8 + (i // 3) * 2.05, 3.8, 1.55, t, b, [BLUE, GREEN, CYAN, AMBER, RED, SLATE][i])
    add_footer(s, 7)

    # 8 call automation
    s = blank()
    bg(s)
    add_title(s, "AI Call Agents", "Inbound and outbound voice automation in many languages.", "The same agent stack can become a sales, support and operations assistant.")
    add_card(s, 0.85, 2.1, 3.85, 2.25, "Inbound calls", "Answer FAQs, qualify leads, capture complaints, route urgent issues, book meetings and update CRM or Sheets.", BLUE)
    add_card(s, 4.95, 2.1, 3.85, 2.25, "Outbound calls", "Lead follow-ups, appointment reminders, surveys, event invitations, payment nudges and campaign outreach.", GREEN)
    add_card(s, 9.05, 2.1, 3.4, 2.25, "Languages", "Telugu, Tamil, Hindi, Kannada, Marathi, English and more based on target markets.", CYAN)
    add_bullets(
        s,
        [
            "Can connect to WhatsApp Business, phone systems, CRM, Google Sheets and dashboards.",
            "Escalates to a human when the customer is angry, high-value or legally sensitive.",
            "Useful for political campaigns, agencies, clinics, real estate, education, finance and local businesses.",
        ],
        1.0,
        5.0,
        11.2,
        1.0,
        13,
        INK,
    )
    add_footer(s, 8)

    # 9 why us
    s = blank()
    bg(s)
    add_title(s, "Why APP.AI", "We combine analytics, agents and automation into one practical delivery model.")
    add_card(s, 0.9, 2.0, 3.55, 2.7, "Proof-first demos", f"Real working dashboard already processes {total} YouTube items in the latest 24-hour demo run.", BLUE)
    add_card(s, 4.9, 2.0, 3.55, 2.7, "Regional-market strength", "Built for Indian languages, film/political ecosystems, PR workflows and local decision makers.", GREEN)
    add_card(s, 8.9, 2.0, 3.55, 2.7, "Actionable outputs", "Clients get links, summaries, risk scores and next actions, not only beautiful charts.", AMBER)
    add_textbox(s, "Our position: build quickly, prove with client data, then automate the recurring work.", 1.0, 5.55, 10.8, 0.45, 20, True, INK)
    add_footer(s, 9)

    # 10 target audience
    s = blank()
    bg(s)
    add_title(s, "Target Clients", "Best first customers are teams who already spend on reputation, calls or growth.")
    add_bullets(
        s,
        [
            "PR agencies and celebrity management firms in Andhra Pradesh, Telangana, Karnataka and Mumbai.",
            "Political consultants, campaign managers, digital war rooms and constituency teams.",
            "Film production houses, music labels, trailer-launch teams and entertainment marketing agencies.",
            "Founder-led brands, real estate groups, education institutions, healthcare networks and local enterprises.",
            "Decision makers: founder, co-founder, managing director, PR head, campaign manager, digital head, celebrity manager and business development head.",
        ],
        0.95,
        2.0,
        7.2,
        3.2,
        15,
        INK,
    )
    add_card(s, 8.7, 2.1, 3.45, 2.7, "Entry offer", "7-day pilot: monitor one person, brand or campaign. Deliver daily Sheet, dashboard and risk-link report.", BLUE)
    add_card(s, 8.7, 5.05, 3.45, 1.0, "Expansion", "Add Meta, X, web, AI calls and sales automation after proof.", GREEN)
    add_footer(s, 10)

    # 11 demo links
    s = blank()
    bg(s, INK)
    add_textbox(s, "Demo Access", 0.8, 0.72, 4.4, 0.55, 30, True, LIGHT)
    add_textbox(s, "Use these during a client meeting to show live proof, raw data and the polished dashboard.", 0.83, 1.35, 8.8, 0.35, 13, False, RGBColor(203, 213, 225))
    add_card(s, 0.95, 2.25, 3.7, 2.75, "Client dashboard", DASHBOARD_LINK, CYAN)
    s.shapes.add_picture(str(qr_dashboard), Inches(1.55), Inches(3.15), width=Inches(1.4))
    add_link(s, "Open client dashboard", DASHBOARD_LINK, 1.0, 4.65, 2.25, 0.28, CYAN)
    add_card(s, 4.95, 2.25, 3.7, 2.75, "Google Sheets data", SHEETS_LINK, GREEN)
    s.shapes.add_picture(str(qr_sheet), Inches(5.55), Inches(3.15), width=Inches(1.4))
    add_link(s, "Open Google Sheet", SHEETS_LINK, 5.0, 4.65, 2.2, 0.28, GREEN)
    add_card(s, 8.95, 2.25, 3.35, 2.75, "Analyst dashboard", ANALYST_LINK, BLUE)
    add_link(s, "Open analyst dashboard", ANALYST_LINK, 9.0, 4.65, 2.35, 0.28, CYAN)
    add_textbox(s, "Recommended client flow: open dashboard first, then show Google Sheets as raw proof.", 1.0, 5.9, 10.3, 0.35, 16, True, MINT)
    add_footer(s, 11)

    # 12 close
    s = blank()
    bg(s, INK)
    s.shapes.add_picture(str(BRAND / "app-ai-icon.png"), Inches(0.8), Inches(0.72), width=Inches(1.1))
    add_textbox(s, "APP.AI", 2.05, 0.86, 2.5, 0.45, 24, True, LIGHT)
    add_textbox(s, "A practical AI partner for data analysis, ML models, agents and automation.", 0.85, 2.05, 9.6, 0.7, 31, True, LIGHT)
    add_bullets(
        s,
        [
            "Start with a live monitoring pilot.",
            "Prove value through links, summaries and daily decisions.",
            "Expand into full-channel intelligence and multilingual AI agents.",
        ],
        1.0,
        3.25,
        8.1,
        1.35,
        18,
        RGBColor(226, 232, 240),
    )
    add_label(s, "READY FOR CLIENT DEMO", 0.96, 5.4, 2.55, 0.38, RGBColor(17, 36, 67), CYAN)
    add_footer(s, 12)

    out = OUTPUT / "APP_AI_Client_Presentation.pptx"
    prs.save(out)
    return out


if __name__ == "__main__":
    out = build_deck()
    print(out)
